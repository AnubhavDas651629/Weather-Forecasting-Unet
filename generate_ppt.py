from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

def create_presentation():
    prs = Presentation()
    
    # slide layouts: 
    # 0: Title Slide
    # 1: Title and Content
    # 2: Section Header
    # 5: Title Only
    # 6: Blank

    def add_section_slide(title_text):
        slide = prs.slides.add_slide(prs.slide_layouts[2])
        slide.shapes.title.text = title_text
        for run in slide.shapes.title.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(54)
            run.font.name = "Helvetica Neue"
        # Remove subtitle if present
        if len(slide.placeholders) > 1:
            sp = slide.placeholders[1]
            sp.element.getparent().remove(sp.element)
        return slide

    def add_content_slide(title_text, content_blocks):
        slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title only
        title = slide.shapes.title
        title.text = title_text
        for run in title.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(36)
            run.font.name = "Helvetica Neue"
            
        # Add textbox for content
        left = Inches(0.5)
        top = Inches(1.3)
        width = Inches(9)
        height = Inches(5.8)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        for i, (header, text) in enumerate(content_blocks):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.space_after = Pt(10)
            
            if header:
                run1 = p.add_run()
                run1.text = header + "\n"
                run1.font.bold = True
                run1.font.italic = True
                run1.font.size = Pt(14)
                run1.font.name = "Helvetica Neue"
                
            if text:
                run2 = p.add_run()
                run2.text = text
                run2.font.size = Pt(14)
                run2.font.name = "Helvetica Neue"
                
    # ---------------------------------------------------------
    # 1. Main Title Slide
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Weather Forecasting U-Net"
    subtitle = slide.placeholders[1]
    subtitle.text = "Comprehensive Codebase Walkthrough\nArchitecture and Helper Functions"
    
    # ---------------------------------------------------------
    # 1.5 BACKGROUND & PAPER OVERVIEW
    # ---------------------------------------------------------
    add_section_slide("PROJECT BACKGROUND")
    
    add_content_slide("THE STANFORD PAPER:", [
        ("Title:", "Weather Forecasting UNET (Meza, Anderson, Tsultrim, 2024 - Stanford CS231n)"),
        ("Problem Statement:", "Traditional severe-weather machine learning often treats each grid point independently (e.g., logistic regression). Tornado risk, however, is inherently spatial. Weather fronts and synoptic patterns extend across massive regions."),
        ("The Solution:", "A U-Net captures these spatial dependencies via convolutions and skip connections, allowing it to \"see\" the surrounding regional weather patterns rather than just isolated pixels."),
        ("Input Features:", "3-Channel Geospatial Grids (256x256): CAPE (Convective Available Potential Energy), CIN (Convective Inhibition), and Tropopause Geopotential Height."),
        ("Target Outputs:", "2-Channel Probability Maps (256x256): Probability of ANY tornado within ~25 miles, and probability of a SIGNIFICANT (EF2+) tornado.")
    ])

    add_content_slide("PROJECT DELIVERABLES:", [
        ("1. Data Preprocessing Pipeline:", "A robust script (preprocess.py) to ingest, synchronize, and resize massive global NetCDF weather files into aligned NumPy tensors."),
        ("2. Custom PyTorch Dataset:", "An optimized DataLoader interface (NOAA_dataset.py) that performs on-the-fly normalization and tensor stacking."),
        ("3. U-Net Neural Network:", "A fully functional PyTorch U-Net architecture implemented from scratch (unet.py, unet_parts.py)."),
        ("4. Custom Training Loop:", "A training pipeline (main.py) featuring Adam optimization, heavily-weighted BCE Loss for extreme class imbalance, and KL Divergence validation metrics."),
        ("5. Inference & Visualization:", "A complete evaluation suite (inference.py) that overlays network predictions onto geographic maps using official NOAA SPC colormaps.")
    ])
    
    # ---------------------------------------------------------
    # 2. OVERALL ARCHITECTURE
    # ---------------------------------------------------------
    add_section_slide("PROJECT ARCHITECTURE")
    
    add_content_slide("CODEBASE OVERVIEW:", [
        ("", "The project is a PyTorch implementation of a U-Net for spatial tornado probability forecasting over the contiguous United States (CONUS)."),
        ("1. preprocess.py:", "Converts raw NetCDF NOAA data (CAPE, CIN, HGT) and Gensini target probabilities into aligned 256x256 NumPy arrays."),
        ("2. NOAA_dataset.py:", "A custom PyTorch Dataset that loads the preprocessed numpy arrays, handles missing values, and normalizes the data on the fly."),
        ("3. unet_parts.py & unet.py:", "Defines the U-Net neural network architecture using reusable blocks (DoubleConv, DownSample, UpSample)."),
        ("4. main.py:", "Executes the core training and validation loops, utilizing Adam optimizer, Weighted BCE Loss, and KL Divergence for evaluation."),
        ("5. inference.py:", "Loads the trained model weights and generates visual probability heatmaps comparing predictions against ground truth.")
    ])

    # ---------------------------------------------------------
    # 3. PREPROCESS.PY
    # ---------------------------------------------------------
    add_section_slide("PREPROCESS.PY")
    
    add_content_slide("ARCHITECTURE:", [
        ("Phase 1: Preparation", "Initializes directories. Uses xarray to open 5 massive raw .nc files, forces all time formats to standard datetime64, and isolates the year 2014. Auto-detects variable names inside the datasets."),
        ("Phase 2: Synchronization", "Extracts the lists of days from all 5 files and cross-references them using np.intersect1d. Throws away any day that doesn't exist perfectly across all 5 files to ensure the neural network never crashes due to missing data."),
        ("Phase 3: Day-by-Day Processing Loop", "Loops through every aligned day. Slices out the flat 2D grids, cleans missing values (NaN to 0.0), resizes grids of different sizes to exactly 256x256 using PyTorch bilinear interpolation, and saves them as .npy files."),
        ("Phase 4: Completion", "Prints a final report detailing exactly how many days succeeded and detailing errors for skipped days.")
    ])
    
    add_content_slide("HELPER FUNCTION", [
        ("def create_folders():", "What it does: Scaffolds the required directory structure on your hard drive.\nWhy we need it: Ensures folders exist before saving thousands of files to prevent crashes."),
        ("def resize_grid(matrix_2d):", "What it does: Cleans NaN values and uses PyTorch bilinear interpolation to squish/stretch maps to 256x256.\nWhy we need it: The U-Net strictly requires 256x256 inputs."),
        ("def get_variable(ds):", "What it does: Scans internal NetCDF metadata to auto-detect the core data array name.\nWhy we need it: Prevents manual hardcoding of variable names for different NOAA datasets."),
        ("def extract_2d(da, date):", "What it does: Slices out a specific date and strips away unnecessary altitude or ensemble dimensions.\nWhy we need it: U-Net requires flat 2D maps, but raw files often contain 3D/4D stacked data."),
        ("def save_validated(arr, path, label):", "What it does: Resizes the map, strictly verifies the shape is 256x256, and saves as .npy.\nWhy we need it: Final safety checkpoint to prevent saving broken files that would ruin training.")
    ])

    # ---------------------------------------------------------
    # 4. NOAA_DATASET.PY
    # ---------------------------------------------------------
    add_section_slide("NOAA_DATASET.PY")
    
    add_content_slide("ARCHITECTURE & FUNCTIONS:", [
        ("def __init__(self, root_path, test=False, file_list=None):", "What it does: Initializes the dataset by constructing absolute file paths to all 5 data channels (cape, cin, geo, tor, sigtor). It performs strict assertion checks to ensure all channels have the exact same number of files.\nWhy we need it: Prepares the dataset state and prevents training on misaligned or missing data."),
        ("def __len__(self):", "What it does: Returns the total number of files in the dataset.\nWhy we need it: Required by PyTorch DataLoader to know how many batches it can create."),
        ("def __getitem__(self, index):", "What it does: Loads the .npy files for a given index. Replaces NaN/inf values with 0.0 using np.nan_to_num. Stacks inputs into a [3, 256, 256] tensor and targets into a [2, 256, 256] tensor. Performs per-channel absolute-max normalization to scale features between [-1.0, 1.0]. Clamps target labels to [0, 1].\nWhy we need it: This is the core pipeline that feeds clean, normalized, ready-to-train tensors into the neural network during every step of training.")
    ])

    # ---------------------------------------------------------
    # 5. UNET_PARTS.PY & UNET.PY
    # ---------------------------------------------------------
    add_section_slide("UNET.PY & UNET_PARTS.PY")
    
    add_content_slide("ARCHITECTURE & HELPER CLASSES:", [
        ("class DoubleConv(nn.Module):", "What it does: Two sequential 3x3 convolutions, each followed by BatchNorm2d and a ReLU activation.\nWhy we need it: The foundational feature-extraction block of the network."),
        ("class DownSample(nn.Module):", "What it does: Applies a DoubleConv and then a 2x2 MaxPool2d to halve the spatial dimensions.\nWhy we need it: Used in the encoder to extract deep, abstract meteorological features while reducing computational load."),
        ("class UpSample(nn.Module):", "What it does: Applies a ConvTranspose2d to double spatial dimensions, calculates spatial differences, pads the upsampled tensor to perfectly match the encoder's skip connection, concatenates them, and applies a DoubleConv.\nWhy we need it: Used in the decoder to restore high-resolution spatial boundaries and merge them with abstract features."),
        ("class UNet(nn.Module):", "What it does: The master assembly. Chains 4 DownSample blocks into a 1024-channel bottleneck, followed by 4 UpSample blocks. Ends with a 1x1 convolution mapping to 2 target classes (logits).\nWhy we need it: This is the complete neural network that maps raw weather inputs to tornado probability maps.")
    ])

    # ---------------------------------------------------------
    # 6. MAIN.PY
    # ---------------------------------------------------------
    add_section_slide("MAIN.PY")
    
    add_content_slide("ARCHITECTURE:", [
        ("Phase 1: Setup & Data Splitting", "Sets hyperparameters (BATCH_SIZE=8, EPOCHS=50). Performs a strict chronological 80/20 split (Jan-Aug for Train, Sep-Oct for Val) to perfectly prevent future data leakage. Initializes DataLoaders."),
        ("Phase 2: Model & Loss Initialization", "Loads the U-Net onto the GPU. Initializes Adam optimizer. Sets up BCEWithLogitsLoss with a pos_weight of [5.0, 10.0] to heavily penalize missing extremely rare tornado pixels."),
        ("Phase 3: The Training Loop", "For each batch, the model makes a guess (logits). It calculates KL Divergence loss. loss.backward() computes gradients, and optimizer.step() updates weights to make the model smarter."),
        ("Phase 4: The Validation Loop", "The model is set to eval() mode (no gradients). It predicts on unseen validation days. Calculates and reports BCE, KL Divergence, and Mean Absolute Error (MAE)."),
        ("Phase 5: Save Model", "Saves the trained neural network weights to models/unet.pth.")
    ])
    
    add_content_slide("HELPER FUNCTION", [
        ("def bernoulli_entropy(p, eps=1e-7):", "What it does: Calculates the statistical entropy of the target labels.\nWhy we need it: Required as a component to mathematically calculate KL Divergence."),
        ("def kl_divergence_from_logits(logits, target, eps=1e-7):", "What it does: Calculates KL Divergence using the mathematical identity KL = BCE - Entropy. Uses raw logits to leverage the log-sum-exp trick.\nWhy we need it: Provides a mathematically stable, crash-proof loss function for backpropagation training."),
        ("def kl_divergence_from_probs(pred_probs, target, eps=1e-7):", "What it does: Calculates KL Divergence explicitly using sigmoid probabilities, clamping values to prevent log(0) crashes.\nWhy we need it: Used strictly for human-readable evaluation and reporting during validation.")
    ])

    # ---------------------------------------------------------
    # 7. INFERENCE.PY
    # ---------------------------------------------------------
    add_section_slide("INFERENCE.PY")
    
    add_content_slide("ARCHITECTURE:", [
        ("Phase 1: Setup & Model Loading", "Defines paths. Instantiates a blank U-Net and loads the trained unet.pth weights into it. Sets the model to evaluation mode (model.eval())."),
        ("Phase 2: Geographic Alignment", "Runs load_coordinates() to fetch exact Latitude/Longitude grids for all 256x256 pixels, enabling Cartopy to accurately draw state borders."),
        ("Phase 3: Prediction", "Loads specific days from the NOAATornadoDataset. Passes them through the model. Applies torch.sigmoid() to convert raw logits into strictly bounded 0.0 to 1.0 probability percentages."),
        ("Phase 4: Visualization", "Applies official NOAA Storm Prediction Center (SPC) custom colormaps. Uses matplotlib and cartopy to render multi-panel side-by-side heatmaps (Inputs vs Predictions vs Ground Truth). Saves final images.")
    ])
    
    add_content_slide("HELPER FUNCTION", [
        ("def load_coordinates():", "What it does: Extracts lat/lon from raw .nc files, wraps longitudes to -180 to 180, and resizes to 256x256.\nWhy we need it: Allows plotting of accurate geographic state borders over the heatmaps."),
        ("def _prepare_input(x): & def _to_numpy_chw(tensor):", "What they do: Safely add/remove artificial 'batch' dimensions from tensors, and move them between GPU memory and CPU NumPy arrays.\nWhy we need them: Bridges the strict shape requirements of PyTorch with the plotting requirements of matplotlib."),
        ("def single_sample_inference(...) & def single_day_inference_from_npy(...):", "What they do: End-to-end visualization wrappers (one by index, one by specific filename).\nWhy we need them: Provides targeted, deep-dive visual evaluations of specific historic storm events.")
    ])

    # ---------------------------------------------------------
    # 8. RESULTS (IMAGES)
    # ---------------------------------------------------------
    add_section_slide("INFERENCE RESULTS")
    
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    title_tf = title_box.text_frame
    title_tf.text = "Visual Outputs: Probability Heatmaps"
    for paragraph in title_tf.paragraphs:
        for run in paragraph.runs:
            run.font.bold = True
            run.font.size = Pt(36)
            run.font.name = "Helvetica Neue"
            run.font.color.rgb = RGBColor(0, 51, 102)
            
    img1_path = "single_sample_0.png"
    img2_path = "grid_output.png"
    
    if os.path.exists(img1_path):
        slide.shapes.add_picture(img1_path, Inches(0.5), Inches(1.2), width=Inches(4.0))
        
    if os.path.exists(img2_path):
        slide.shapes.add_picture(img2_path, Inches(4.8), Inches(1.2), width=Inches(5.0))

    prs.save('Weather_Forecasting_UNet.pptx')
    print("Detailed Presentation saved successfully.")

if __name__ == '__main__':
    create_presentation()
